"""Generate Food App analysis PowerPoint presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
DARK_GREEN = RGBColor(0x1B, 0x5E, 0x20)
MED_GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0x4C, 0xAF, 0x50)
ACCENT_GREEN = RGBColor(0x81, 0xC7, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x21, 0x21, 0x21)
MED_TEXT = RGBColor(0x42, 0x42, 0x42)
LIGHT_TEXT = RGBColor(0x75, 0x75, 0x75)
ORANGE = RGBColor(0xFF, 0x6F, 0x00)
RED_ACCENT = RGBColor(0xC6, 0x28, 0x28)
BLUE = RGBColor(0x15, 0x65, 0xC0)
CARD_BG = RGBColor(0xE8, 0xF5, 0xE9)
SECTION_BG = RGBColor(0xF1, 0xF8, 0xE9)
TOTAL_SLIDES = 16

def add_shape(slide, left, top, width, height, fill_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    return shape

def add_rounded(slide, left, top, w, h, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def tb(slide, left, top, w, h, text, sz=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    return txBox

def add_para(tf, text, sz=14, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, sb=Pt(4)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Calibri'
    p.alignment = align
    p.space_before = sb
    p.space_after = Pt(2)
    return p

def header(slide, title):
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.0), fill_color=DARK_GREEN)
    tb(slide, Inches(0.6), Inches(0.15), Inches(12), Inches(0.7), title, sz=28, color=WHITE, bold=True)

def footer(slide, num):
    tb(slide, Inches(0.5), Inches(7.0), Inches(8), Inches(0.4),
       "Food Waste Reduction Marketplace  |  SMU MBA ML Individual Project", sz=9, color=LIGHT_TEXT)
    tb(slide, Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4),
       f"{num} / {TOTAL_SLIDES}", sz=9, color=LIGHT_TEXT, align=PP_ALIGN.RIGHT)

def card(slide, left, top, w, h, title_text, bullets, title_color=DARK_GREEN):
    add_rounded(slide, left, top, w, h, CARD_BG)
    tb(slide, left + Inches(0.15), top + Inches(0.1), w - Inches(0.3), Inches(0.4),
       title_text, sz=14, color=title_color, bold=True)
    txb = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.5), w - Inches(0.3), h - Inches(0.6))
    tf = txb.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = ""
    for b in bullets:
        add_para(tf, f"• {b}", sz=11, color=MED_TEXT)


# ============================================================
# SLIDE 1: TITLE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill_color=DARK_GREEN)
add_shape(sl, Inches(0), Inches(5.5), Inches(13.333), Inches(2), fill_color=RGBColor(0x14, 0x4D, 0x1A))

tb(sl, Inches(1), Inches(1.2), Inches(11.3), Inches(1.0),
   "Food Waste Reduction Marketplace", sz=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(2.5), Inches(10.3), Inches(0.8),
   "An ML-Powered Platform Connecting Surplus Food to Budget-Conscious Consumers",
   sz=22, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)
add_shape(sl, Inches(4.5), Inches(3.5), Inches(4.333), Pt(3), fill_color=LIGHT_GREEN)
tb(sl, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.5),
   "SMU MBA Machine Learning  |  Individual Project  |  Startup Exploration",
   sz=16, color=RGBColor(0xA5, 0xD6, 0xA7), align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(5.7), Inches(10.3), Inches(0.5),
   "Predict   |   Personalize   |   Price   |   Prevent",
   sz=22, color=ACCENT_GREEN, bold=True, align=PP_ALIGN.CENTER)
tb(sl, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.4),
   "April 2026", sz=14, color=RGBColor(0xA5, 0xD6, 0xA7), align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: AGENDA
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Agenda")
footer(sl, 2)

items = [
    ("01", "Problem & Vision", "Why food waste matters in Singapore"),
    ("02", "Market Opportunity", "TAM/SAM/SOM and competitive landscape"),
    ("03", "Our Solution", "ML-powered marketplace platform"),
    ("04", "Value Propositions", "What makes us different"),
    ("05", "Platform Model", "Two-sided marketplace with network effects"),
    ("06", "ML Architecture", "4 ML capabilities in detail"),
    ("07", "User Experience", "Merchant and consumer journeys"),
    ("08", "Business Model", "Revenue, unit economics, roadmap"),
    ("09", "Risks & Mitigations", "Honest assessment of challenges"),
    ("10", "Demo Vision", "What the prototype will demonstrate"),
]
for i, (num, title, desc) in enumerate(items):
    row = i // 2
    col = i % 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(1.3 + row * 1.1)
    add_rounded(sl, x, y, Inches(5.8), Inches(0.95), CARD_BG)
    tb(sl, x + Inches(0.15), y + Inches(0.08), Inches(0.5), Inches(0.4), num, sz=16, color=MED_GREEN, bold=True)
    tb(sl, x + Inches(0.7), y + Inches(0.05), Inches(4.8), Inches(0.4), title, sz=15, color=DARK_GREEN, bold=True)
    tb(sl, x + Inches(0.7), y + Inches(0.45), Inches(4.8), Inches(0.4), desc, sz=11, color=MED_TEXT)


# ============================================================
# SLIDE 3: THE PROBLEM
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "The Problem: Food Waste in Singapore")
footer(sl, 3)

# Key stat callouts
for i, (val, label) in enumerate([
    ("755,000", "tonnes of food waste\nannually (2023)"),
    ("18%", "recycling rate —\n82% goes to incineration"),
    ("~226,000", "tonnes from F&B &\nretail (our addressable supply)"),
    ("S$1.1B", "estimated value of\nwasted food annually"),
]):
    x = Inches(0.6 + i * 3.15)
    add_rounded(sl, x, Inches(1.3), Inches(2.9), Inches(1.8), CARD_BG)
    tb(sl, x + Inches(0.1), Inches(1.4), Inches(2.7), Inches(0.6), val, sz=32, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(2.1), Inches(2.7), Inches(0.8), label, sz=12, color=MED_TEXT, align=PP_ALIGN.CENTER)

# Why it matters
tb(sl, Inches(0.6), Inches(3.4), Inches(12), Inches(0.4), "Why This Matters", sz=18, color=DARK_GREEN, bold=True)

card(sl, Inches(0.6), Inches(3.9), Inches(6), Inches(2.8), "Environmental Impact", [
    "Food waste is 11% of Singapore's total waste",
    "Incineration produces CO2 — Semakau landfill filling up",
    "Singapore Green Plan 2030 targets waste reduction",
    "NEA tightening mandatory food waste reporting",
])

card(sl, Inches(6.9), Inches(3.9), Inches(5.8), Inches(2.8), "Economic Impact", [
    "F&B businesses lose S$200-500/month to disposal costs",
    "Surplus food has value — it's just sold at the wrong time",
    "Budget-conscious consumers underserved by current options",
    "No ML-powered solution exists in Singapore market",
])


# ============================================================
# SLIDE 4: MARKET OPPORTUNITY
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Market Opportunity: Singapore F&B")
footer(sl, 4)

# TAM SAM SOM
for i, (label, val, desc) in enumerate([
    ("TAM", "S$580M", "Total food waste value\n(F&B + retail)"),
    ("SAM", "S$80-120M", "Realistically recoverable\nsurplus food"),
    ("SOM (Yr 3)", "S$5-15M", "Conservative platform\ncapture (GMV)"),
]):
    x = Inches(0.6 + i * 4.1)
    add_rounded(sl, x, Inches(1.3), Inches(3.8), Inches(2.0), CARD_BG)
    tb(sl, x + Inches(0.1), Inches(1.35), Inches(3.6), Inches(0.3), label, sz=13, color=LIGHT_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(1.7), Inches(3.6), Inches(0.6), val, sz=30, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(2.3), Inches(3.6), Inches(0.7), desc, sz=12, color=MED_TEXT, align=PP_ALIGN.CENTER)

# Market drivers
tb(sl, Inches(0.6), Inches(3.6), Inches(12), Inches(0.4), "Key Market Drivers", sz=18, color=DARK_GREEN, bold=True)

for i, (icon, title, desc) in enumerate([
    ("Regulatory", "Policy Tailwind", "NEA mandatory reporting, Green Plan 2030, tightening waste regulations"),
    ("Consumer", "Growing Demand", "65% willing to buy near-expiry at discount, rising cost-of-living pressure"),
    ("Competitive", "Clear Gap", "No ML-powered competitor. TreatSure (only SG player) is small, no ML"),
    ("Technology", "Ready", "96% smartphone penetration, established food delivery app behavior"),
]):
    x = Inches(0.6 + i * 3.15)
    add_rounded(sl, x, Inches(4.1), Inches(2.9), Inches(2.6), CARD_BG)
    tb(sl, x + Inches(0.1), Inches(4.15), Inches(2.7), Inches(0.3), icon, sz=10, color=LIGHT_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(4.45), Inches(2.7), Inches(0.35), title, sz=14, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(4.85), Inches(2.7), Inches(1.5), desc, sz=11, color=MED_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 5: COMPETITIVE LANDSCAPE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Competitive Landscape")
footer(sl, 5)

# Competitor cards
competitors = [
    ("Too Good To Go", "Global leader, 75M+ users\n17 countries, Surprise Bags", "No ML, no SG presence,\nmanual listing & pricing"),
    ("TreatSure", "Singapore-based, only\ndirect competitor", "Small (est. 50-100K users),\nno ML, buffet-focused only"),
    ("OLIO", "Community-driven food\nsharing, 7M+ users", "Not merchant-focused,\nno pricing optimization"),
    ("Karma / Phenix", "European platforms with\ngrocery partnerships", "Not in Asia, basic analytics,\nno personalization"),
]
for i, (name, strength, gap) in enumerate(competitors):
    x = Inches(0.5 + i * 3.15)
    add_rounded(sl, x, Inches(1.3), Inches(2.95), Inches(2.5), CARD_BG)
    tb(sl, x + Inches(0.1), Inches(1.35), Inches(2.75), Inches(0.35), name, sz=16, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(1.75), Inches(2.75), Inches(0.3), "Strengths:", sz=10, color=LIGHT_GREEN, bold=True)
    tb(sl, x + Inches(0.1), Inches(1.95), Inches(2.75), Inches(0.6), strength, sz=10, color=MED_TEXT)
    tb(sl, x + Inches(0.1), Inches(2.55), Inches(2.75), Inches(0.3), "Gaps:", sz=10, color=RED_ACCENT, bold=True)
    tb(sl, x + Inches(0.1), Inches(2.75), Inches(2.75), Inches(0.6), gap, sz=10, color=MED_TEXT)

# Our differentiator
add_rounded(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.7), RGBColor(0xE3, 0xF2, 0xFD))
tb(sl, Inches(0.7), Inches(4.2), Inches(11.9), Inches(0.4), "Our Differentiation: ML-First Approach", sz=18, color=BLUE, bold=True)

diffs = [
    ("Surplus Prediction", "Predict surplus before it happens — auto-generate listings, one-tap approval"),
    ("Dynamic Pricing", "Time/demand-based pricing optimization — higher sell-through than fixed discounts"),
    ("Personalization", "ML-ranked deal feed based on preferences, behavior, and context"),
    ("Prescriptive Analytics", "\"Reduce batch by 20%\" not \"You wasted X kg\" — actionable merchant insights"),
]
for i, (title, desc) in enumerate(diffs):
    x = Inches(0.7 + i * 3.1)
    tb(sl, x, Inches(4.7), Inches(2.9), Inches(0.35), title, sz=13, color=BLUE, bold=True)
    tb(sl, x, Inches(5.05), Inches(2.9), Inches(1.5), desc, sz=11, color=MED_TEXT)


# ============================================================
# SLIDE 6: OUR SOLUTION
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Our Solution: ML-Powered Food Waste Marketplace")
footer(sl, 6)

# Platform overview
tb(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.8),
   "A two-sided marketplace connecting food businesses with surplus food to budget-conscious consumers, "
   "powered by four ML capabilities that create a self-improving data moat.",
   sz=14, color=MED_TEXT)

# Flow diagram (simplified)
flow_items = [
    ("Food Business", "Has surplus food\n(end-of-day, near-expiry, overstock)"),
    ("Platform", "Predicts surplus, prices\ndynamically, recommends deals"),
    ("Consumer", "Discovers deals,\npurchases at 50-70% off"),
]
for i, (title, desc) in enumerate(flow_items):
    x = Inches(1.2 + i * 4.0)
    add_rounded(sl, x, Inches(2.2), Inches(3.2), Inches(1.8), CARD_BG if i != 1 else RGBColor(0xC8, 0xE6, 0xC9))
    tb(sl, x + Inches(0.1), Inches(2.25), Inches(3.0), Inches(0.4), title, sz=16, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(2.7), Inches(3.0), Inches(1.0), desc, sz=12, color=MED_TEXT, align=PP_ALIGN.CENTER)

# Arrows between
for x_pos in [Inches(4.4), Inches(8.4)]:
    tb(sl, x_pos, Inches(2.8), Inches(0.8), Inches(0.4), "→", sz=30, color=MED_GREEN, bold=True, align=PP_ALIGN.CENTER)

# 4 ML pillars
tb(sl, Inches(0.6), Inches(4.3), Inches(12), Inches(0.4), "Four ML Pillars", sz=18, color=DARK_GREEN, bold=True)

pillars = [
    ("Predict", "Surplus Prediction", "XGBoost + Prophet ensemble\npredicts daily surplus per\nmerchant per category", "80%+ directional accuracy"),
    ("Personalize", "Recommendations", "Hybrid engine combining\ncontent-based + collaborative\nfiltering + context signals", ">15% click-through rate"),
    ("Price", "Dynamic Pricing", "Time-decay curve with demand\nmodifiers, merchant price\nfloors respected", ">80% sell-through rate"),
    ("Prevent", "Waste Analytics", "Pattern detection, anomaly\nflagging, prescriptive\nrecommendations", "Actionable insights weekly"),
]
colors = [BLUE, MED_GREEN, ORANGE, RGBColor(0x7B, 0x1F, 0xA2)]
for i, (icon, title, desc, metric) in enumerate(pillars):
    x = Inches(0.6 + i * 3.15)
    add_rounded(sl, x, Inches(4.8), Inches(2.9), Inches(2.0), CARD_BG)
    tb(sl, x + Inches(0.1), Inches(4.85), Inches(2.7), Inches(0.35), f"{icon}: {title}", sz=13, color=colors[i], bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(5.25), Inches(2.7), Inches(1.0), desc, sz=11, color=MED_TEXT, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(6.2), Inches(2.7), Inches(0.35), f"Target: {metric}", sz=10, color=colors[i], bold=True, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 7: VALUE PROPOSITIONS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Value Propositions")
footer(sl, 7)

# Merchant VPs
tb(sl, Inches(0.6), Inches(1.2), Inches(6), Inches(0.4), "For Merchants (Supply Side)", sz=18, color=DARK_GREEN, bold=True)

mvps = [
    ("Turn Waste into Revenue", "Recover S$50-200/month from food that would be thrown away"),
    ("Zero-Effort Management", "AI predicts surplus, suggests price, generates listing — one tap to approve"),
    ("Actionable Waste Insights", "\"Reduce pastry batch 20% on Tuesdays\" — not just \"you wasted X kg\""),
]
for i, (title, desc) in enumerate(mvps):
    y = Inches(1.7 + i * 1.15)
    add_rounded(sl, Inches(0.6), y, Inches(5.8), Inches(1.0), CARD_BG)
    tb(sl, Inches(0.75), y + Inches(0.05), Inches(5.5), Inches(0.35), title, sz=14, color=DARK_GREEN, bold=True)
    tb(sl, Inches(0.75), y + Inches(0.4), Inches(5.5), Inches(0.5), desc, sz=12, color=MED_TEXT)

# Consumer VPs
tb(sl, Inches(6.9), Inches(1.2), Inches(6), Inches(0.4), "For Consumers (Demand Side)", sz=18, color=BLUE, bold=True)

cvps = [
    ("Premium Food, Fraction of Price", "Quality meals from restaurants you love at 50-70% off"),
    ("Deals That Find You", "ML learns your preferences and surfaces the right deal at the right time"),
    ("Make a Tangible Impact", "Every meal rescued prevents CO2 emissions — track your personal impact"),
]
for i, (title, desc) in enumerate(cvps):
    y = Inches(1.7 + i * 1.15)
    add_rounded(sl, Inches(6.9), y, Inches(5.8), Inches(1.0), RGBColor(0xE3, 0xF2, 0xFD))
    tb(sl, Inches(7.05), y + Inches(0.05), Inches(5.5), Inches(0.35), title, sz=14, color=BLUE, bold=True)
    tb(sl, Inches(7.05), y + Inches(0.4), Inches(5.5), Inches(0.5), desc, sz=12, color=MED_TEXT)

# USP Summary
add_rounded(sl, Inches(0.6), Inches(5.3), Inches(12.1), Inches(1.5), RGBColor(0xFD, 0xF2, 0xE3))
tb(sl, Inches(0.8), Inches(5.35), Inches(11.7), Inches(0.35),
   "Defensible USPs (Hard to Replicate)", sz=14, color=ORANGE, bold=True)
tb(sl, Inches(0.8), Inches(5.75), Inches(11.7), Inches(0.9),
   "1. ML-predicted surplus listings — no competitor offers this; requires merchant data + ML accuracy\n"
   "2. Prescriptive waste analytics — goes beyond dashboards to specific, actionable recommendations\n"
   "3. Data moat from transactions — compounding advantage: more data → better predictions → more merchants",
   sz=12, color=MED_TEXT)


# ============================================================
# SLIDE 8: PLATFORM MODEL
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Platform Model & Network Effects")
footer(sl, 8)

# AAA Framework
tb(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4), "AAA Framework", sz=18, color=DARK_GREEN, bold=True)

for i, (letter, word, subtitle, desc) in enumerate([
    ("A", "Automate", "Reduce operational costs",
     "Auto-generate listings from ML predictions\nDynamic pricing eliminates manual price decisions\nMerchant effort: <5 min/day vs 30+ min on manual platforms"),
    ("A", "Augment", "Reduce decision-making costs",
     "Personalized deal rankings for consumers\nSurplus predictions and pricing guidance for merchants\nWaste pattern insights for operational decisions"),
    ("A", "Amplify", "Reduce expertise costs",
     "Merchant gets data-team-level analytics without data team\nDynamic pricing without pricing expertise\nImpact quantification without sustainability consulting"),
]):
    x = Inches(0.6 + i * 4.1)
    add_rounded(sl, x, Inches(1.7), Inches(3.8), Inches(2.6), CARD_BG)
    tb(sl, x + Inches(0.1), Inches(1.75), Inches(3.6), Inches(0.4), f"{letter} — {word}", sz=16, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(2.15), Inches(3.6), Inches(0.3), subtitle, sz=11, color=LIGHT_GREEN, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(2.5), Inches(3.6), Inches(1.5), desc, sz=11, color=MED_TEXT, align=PP_ALIGN.LEFT)

# Growth flywheel
tb(sl, Inches(0.6), Inches(4.5), Inches(12), Inches(0.4), "Self-Reinforcing Growth Flywheel", sz=18, color=DARK_GREEN, bold=True)

add_rounded(sl, Inches(0.6), Inches(5.0), Inches(12.1), Inches(1.7), RGBColor(0xE8, 0xF5, 0xE9))
flywheel = (
    "More merchants → More deals → More consumers → More transaction data → "
    "Better ML predictions → Higher sell-through → More merchant revenue → "
    "More merchants →  (compounding advantage)"
)
tb(sl, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.5), flywheel, sz=13, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
tb(sl, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.8),
   "The data moat is self-reinforcing: each prediction cycle improves accuracy. "
   "New entrants start from zero data. This creates a compounding advantage that is expensive to replicate. "
   "Singapore's small market size actually helps — density drives faster data accumulation.",
   sz=12, color=MED_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9: ML SURPLUS PREDICTION
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "ML Capability 1: Surplus Prediction")
footer(sl, 9)

tb(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
   "Predicts how much surplus food a merchant will have for each food category, "
   "enabling auto-generated listings and proactive waste reduction.",
   sz=13, color=MED_TEXT)

# Feature categories
tb(sl, Inches(0.6), Inches(1.8), Inches(6), Inches(0.35), "Input Features", sz=16, color=DARK_GREEN, bold=True)
features = [
    "Temporal: Day of week, holidays, school breaks, payday proximity",
    "Weather: Temperature, rainfall, humidity, weather category",
    "Historical: Sales last 7 days, same day last week, 30-day average",
    "Merchant: Avg surplus %, batch size, operating hours, tenure",
    "Events: Nearby event count, attendance estimates",
    "Lagged surplus: Yesterday, same day last week, 7-day trend",
]
for i, f in enumerate(features):
    tb(sl, Inches(0.8), Inches(2.2 + i * 0.35), Inches(5.5), Inches(0.35), f"• {f}", sz=11, color=MED_TEXT)

# Model architecture
tb(sl, Inches(6.9), Inches(1.8), Inches(6), Inches(0.35), "Model Architecture", sz=16, color=DARK_GREEN, bold=True)

card(sl, Inches(6.9), Inches(2.2), Inches(5.8), Inches(1.5), "Phase 1: Rule-Based (Cold Start)", [
    "Category averages × day-of-week multipliers",
    "Weather modifiers (rainy = 1.2x surplus)",
    "Used for first 2-4 weeks per merchant",
])

card(sl, Inches(6.9), Inches(3.85), Inches(5.8), Inches(1.5), "Phase 2: XGBoost + Prophet Ensemble", [
    "XGBoost (65% weight): Tabular feature handling",
    "Prophet (35% weight): Weekly seasonality + trend",
    "Per-merchant tuning after 4+ weeks of data",
])

# Accuracy targets
tb(sl, Inches(0.6), Inches(4.7), Inches(12), Inches(0.35), "Performance Targets", sz=16, color=DARK_GREEN, bold=True)
add_rounded(sl, Inches(0.6), Inches(5.1), Inches(12.1), Inches(1.5), CARD_BG)
metrics = [
    ("MAE", "<2 units/category", "Average prediction error"),
    ("Directional Accuracy", ">85%", "Correctly predicts surplus vs. no-surplus"),
    ("MAPE", "<25%", "Relative accuracy percentage"),
    ("Coverage", ">95%", "Days with predictions available"),
]
for i, (metric, target, desc) in enumerate(metrics):
    x = Inches(0.8 + i * 3.0)
    tb(sl, x, Inches(5.2), Inches(2.8), Inches(0.3), metric, sz=12, color=DARK_GREEN, bold=True)
    tb(sl, x, Inches(5.5), Inches(2.8), Inches(0.3), target, sz=18, color=MED_GREEN, bold=True)
    tb(sl, x, Inches(5.9), Inches(2.8), Inches(0.3), desc, sz=10, color=MED_TEXT)


# ============================================================
# SLIDE 10: ML RECOMMENDATIONS + PRICING
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "ML Capabilities 2 & 3: Recommendations & Dynamic Pricing")
footer(sl, 10)

# Left half: Recommendations
tb(sl, Inches(0.6), Inches(1.2), Inches(6), Inches(0.35), "Personalized Recommendations", sz=18, color=DARK_GREEN, bold=True)

card(sl, Inches(0.6), Inches(1.65), Inches(5.8), Inches(2.0), "Hybrid Recommendation Engine", [
    "Content-based: Match user preferences to listing features",
    "Collaborative: \"Users like you also liked\" via matrix factorization",
    "Contextual: Time of day, weather, location, urgency signals",
    "Final score = 35% content + 35% collaborative + 30% contextual",
])

card(sl, Inches(0.6), Inches(3.85), Inches(5.8), Inches(2.5), "Cold Start Solutions", [
    "New user: 3-question preference quiz → content-based defaults",
    "New merchant: Similar merchant proxy + exploration boost",
    "After 10 interactions: Collaborative filtering activates",
    "Diversity enforced: Max 3 same-merchant in top 20, min 3 cuisines in top 10",
])

# Right half: Dynamic Pricing
tb(sl, Inches(6.9), Inches(1.2), Inches(6), Inches(0.35), "Dynamic Pricing", sz=18, color=ORANGE, bold=True)

card(sl, Inches(6.9), Inches(1.65), Inches(5.8), Inches(1.8), "Time-Decay Price Curve", [
    "price(t) = base × (floor% + (1-floor%) × e^(-λt))",
    "3 tiers: Initial (60-70%) → Mid (40-50%) → Floor price",
    "λ adjusted by demand signals (views, purchases, sell-through)",
    "Prices only decrease — never increase — preserving trust",
])

card(sl, Inches(6.9), Inches(3.65), Inches(5.8), Inches(1.3), "Merchant Controls", [
    "Price floor per category (never go below)",
    "Pricing mode: Conservative / Balanced / Aggressive",
    "Toggle dynamic pricing on/off per category",
])

card(sl, Inches(6.9), Inches(5.15), Inches(5.8), Inches(1.2), "Target Outcome", [
    "Sell-through rate: >80% (vs ~60% for fixed pricing)",
    "Revenue per item: Higher than flat discount",
    "Consumer trust: Transparent pricing with visible savings %",
])


# ============================================================
# SLIDE 11: ML WASTE ANALYTICS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "ML Capability 4: Prescriptive Waste Analytics")
footer(sl, 11)

tb(sl, Inches(0.6), Inches(1.2), Inches(12), Inches(0.6),
   "Goes beyond descriptive dashboards to provide specific, actionable recommendations that reduce waste at source — "
   "the highest-value capability for merchants.",
   sz=13, color=MED_TEXT)

# 4 analytics capabilities
analytics = [
    ("Waste Breakdown", "Which categories contribute most to waste",
     "Pastries: 42% (38 kg, S$520 lost)\nRice dishes: 28% (25 kg, S$340 lost)\n→ Focus reduction efforts here"),
    ("Pattern Detection", "K-Means clustering identifies merchant waste profile",
     "Profiles: Consistent over-producer,\nWeekend spike, Weather-sensitive\n→ Tailored recommendations"),
    ("Anomaly Detection", "Isolation Forest flags unusual waste events",
     "Apr 12: 28 units waste (expected 12-18)\nCause: Heavy rain + late listing\n→ Proactive alert"),
    ("Prescriptive Recommendations", "Specific actions with estimated savings",
     "\"Reduce pastry batch 20% on weekdays\"\nEstimated savings: S$180/month\nConfidence: 85%"),
]
colors_a = [BLUE, MED_GREEN, ORANGE, RGBColor(0x7B, 0x1F, 0xA2)]
for i, (title, desc, example) in enumerate(analytics):
    x = Inches(0.5 + i * 3.2)
    add_rounded(sl, x, Inches(2.0), Inches(3.0), Inches(3.8), CARD_BG)
    tb(sl, x + Inches(0.1), Inches(2.05), Inches(2.8), Inches(0.4), title, sz=14, color=colors_a[i], bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(2.5), Inches(2.8), Inches(0.6), desc, sz=11, color=MED_TEXT, align=PP_ALIGN.CENTER)
    add_shape(sl, x + Inches(0.3), Inches(3.2), Inches(2.4), Pt(1), fill_color=ACCENT_GREEN)
    tb(sl, x + Inches(0.1), Inches(3.4), Inches(2.8), Inches(2.0), example, sz=11, color=DARK_TEXT, align=PP_ALIGN.LEFT)

# Impact
tb(sl, Inches(0.6), Inches(6.1), Inches(12), Inches(0.6),
   "Key Insight: Prescriptive analytics is more valuable than marketplace sales — "
   "it addresses the root cause and makes the platform valuable even without consumer demand, "
   "solving the marketplace chicken-and-egg problem.",
   sz=12, color=DARK_GREEN, bold=True)


# ============================================================
# SLIDE 12: USER EXPERIENCE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "User Experience: Merchant & Consumer Journeys")
footer(sl, 12)

# Merchant journey
tb(sl, Inches(0.6), Inches(1.2), Inches(6), Inches(0.35), "Merchant Journey", sz=18, color=DARK_GREEN, bold=True)

m_steps = [
    ("1. Onboard", "Business details, SFA license\nverification, category setup\n(<15 minutes)"),
    ("2. Predict", "Morning: AI predicts surplus\nby category with confidence\nscore and suggested price"),
    ("3. Approve", "One-tap to accept prediction\n→ listing goes live with\ndynamic pricing"),
    ("4. Fulfill", "Consumer purchases →\nQR code pickup →\norder complete"),
    ("5. Improve", "Waste analytics dashboard\nwith prescriptive insights\nand trend tracking"),
]
for i, (step, desc) in enumerate(m_steps):
    x = Inches(0.5 + i * 1.28)
    add_rounded(sl, x, Inches(1.65), Inches(1.18), Inches(2.3), CARD_BG)
    tb(sl, x + Inches(0.05), Inches(1.7), Inches(1.08), Inches(0.35), step, sz=10, color=MED_GREEN, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.05), Inches(2.1), Inches(1.08), Inches(1.5), desc, sz=9, color=MED_TEXT, align=PP_ALIGN.CENTER)

# Consumer journey
tb(sl, Inches(0.6), Inches(4.2), Inches(6), Inches(0.35), "Consumer Journey", sz=18, color=BLUE, bold=True)

c_steps = [
    ("1. Sign Up", "Quick preference quiz\n(3 questions, <2 min)\nLocation permission"),
    ("2. Discover", "ML-personalized feed\nFor You, Nearby,\nExpiring Soon sections"),
    ("3. Purchase", "3-tap purchase:\nBrowse → Select → Pay\nQR code generated"),
    ("4. Pick Up", "Show QR at merchant\nScan → collect food\nRate experience"),
    ("5. Impact", "Track meals rescued,\nmoney saved, CO2\nprevented, earn badges"),
]
for i, (step, desc) in enumerate(c_steps):
    x = Inches(0.5 + i * 1.28)
    add_rounded(sl, x, Inches(4.65), Inches(1.18), Inches(2.0), RGBColor(0xE3, 0xF2, 0xFD))
    tb(sl, x + Inches(0.05), Inches(4.7), Inches(1.08), Inches(0.35), step, sz=10, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.05), Inches(5.1), Inches(1.08), Inches(1.3), desc, sz=9, color=MED_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 13: BUSINESS MODEL
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Business Model & Unit Economics")
footer(sl, 13)

# Revenue model
tb(sl, Inches(0.6), Inches(1.2), Inches(6), Inches(0.35), "Revenue Model", sz=18, color=DARK_GREEN, bold=True)

card(sl, Inches(0.6), Inches(1.65), Inches(5.8), Inches(2.3), "Commission-Based Marketplace", [
    "15% commission per transaction (industry standard)",
    "Merchant receives 85% of sale price",
    "Minimum commission: S$0.30 per transaction",
    "No upfront cost for merchants — zero barrier to join",
    "Future: Freemium subscription for advanced analytics",
])

# Unit economics
tb(sl, Inches(6.9), Inches(1.2), Inches(6), Inches(0.35), "Unit Economics (3 Scenarios)", sz=18, color=DARK_GREEN, bold=True)

add_rounded(sl, Inches(6.9), Inches(1.65), Inches(5.8), Inches(2.3), CARD_BG)
tb(sl, Inches(7.0), Inches(1.7), Inches(1.8), Inches(0.3), "Metric", sz=10, color=MED_TEXT, bold=True)
tb(sl, Inches(8.8), Inches(1.7), Inches(1.2), Inches(0.3), "Pessimistic", sz=10, color=RED_ACCENT, bold=True, align=PP_ALIGN.CENTER)
tb(sl, Inches(10.0), Inches(1.7), Inches(1.2), Inches(0.3), "Moderate", sz=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
tb(sl, Inches(11.2), Inches(1.7), Inches(1.2), Inches(0.3), "Optimistic", sz=10, color=MED_GREEN, bold=True, align=PP_ALIGN.CENTER)

for i, (metric, pess, mod, opt) in enumerate([
    ("Avg deal value", "S$3", "S$4", "S$6"),
    ("Txns/merchant/week", "3", "8", "15"),
    ("Merchants for break-even", "370+", "100", "50"),
    ("Consumer CAC", "S$15", "S$10", "S$5"),
]):
    y = Inches(2.05 + i * 0.35)
    tb(sl, Inches(7.0), y, Inches(1.8), Inches(0.35), metric, sz=10, color=MED_TEXT)
    tb(sl, Inches(8.8), y, Inches(1.2), Inches(0.35), pess, sz=10, color=RED_ACCENT, align=PP_ALIGN.CENTER)
    tb(sl, Inches(10.0), y, Inches(1.2), Inches(0.35), mod, sz=10, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, Inches(11.2), y, Inches(1.2), Inches(0.35), opt, sz=10, color=MED_GREEN, align=PP_ALIGN.CENTER)

# Expansion
tb(sl, Inches(0.6), Inches(4.2), Inches(12), Inches(0.35), "Expansion Strategy", sz=18, color=DARK_GREEN, bold=True)

card(sl, Inches(0.6), Inches(4.65), Inches(3.8), Inches(1.8), "Months 1-6: Singapore MVP", [
    "Target: 20-50 merchants (restaurants, cafes, bakeries)",
    "Rule-based features, manual onboarding",
    "Build initial data for ML training",
])

card(sl, Inches(4.7), Inches(4.65), Inches(3.8), Inches(1.8), "Months 6-18: ML Activation", [
    "ML predictions activate as data accumulates",
    "Target: 100-200 merchants",
    "Dynamic pricing and analytics mature",
])

card(sl, Inches(8.8), Inches(4.65), Inches(3.9), Inches(1.8), "18+ Months: Regional", [
    "Expand to Malaysia, Thailand",
    "Hawker stall inclusion",
    "Corporate meal programs",
])


# ============================================================
# SLIDE 14: IMPLEMENTATION ROADMAP
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Implementation Roadmap: 10-Week Plan")
footer(sl, 14)

phases = [
    ("Phase 1: Core Marketplace", "Weeks 1-3", [
        "PostgreSQL + FastAPI backend",
        "Data model + API endpoints",
        "Merchant dashboard (web)",
        "Consumer app (Flutter)",
        "Stripe payments + QR pickup",
    ], MED_GREEN),
    ("Phase 2: ML Features", "Weeks 4-7", [
        "Surplus prediction (XGBoost)",
        "Recommendations (content-based)",
        "Dynamic pricing (rule-based)",
        "Waste analytics (clustering)",
        "Synthetic data for demo",
    ], BLUE),
    ("Phase 3: Polish & Demo", "Weeks 8-10", [
        "ML pipeline infrastructure",
        "UX polish and performance",
        "Model tuning and evaluation",
        "Demo scenario preparation",
        "Presentation materials",
    ], RGBColor(0x7B, 0x1F, 0xA2)),
]

for i, (title, weeks, items, color) in enumerate(phases):
    x = Inches(0.5 + i * 4.2)
    add_rounded(sl, x, Inches(1.3), Inches(3.9), Inches(5.2), CARD_BG)

    # Phase header
    add_shape(sl, x, Inches(1.3), Inches(3.9), Inches(0.8), fill_color=color)
    tb(sl, x + Inches(0.1), Inches(1.35), Inches(3.7), Inches(0.4), title, sz=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), Inches(1.7), Inches(3.7), Inches(0.3), weeks, sz=12, color=RGBColor(0xE0, 0xE0, 0xE0), align=PP_ALIGN.CENTER)

    for j, item in enumerate(items):
        tb(sl, x + Inches(0.15), Inches(2.3 + j * 0.5), Inches(3.6), Inches(0.4),
           f"✓  {item}", sz=12, color=MED_TEXT)

# Tech stack
tb(sl, Inches(0.6), Inches(6.7), Inches(12), Inches(0.3),
   "Tech: Python/FastAPI  |  PostgreSQL  |  Flutter  |  scikit-learn/XGBoost  |  Stripe  |  Docker",
   sz=11, color=LIGHT_TEXT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15: RISKS & MITIGATIONS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
header(sl, "Risks & Mitigations")
footer(sl, 15)

risks = [
    ("Cold Start / No Data", "HIGH", "ML models need data; platform has zero users at launch",
     "Synthetic data for demo; rule-based features initially; ML activates as real data accumulates"),
    ("TGTG Enters Singapore", "MODERATE", "Global leader with 75M users could enter market",
     "First-mover advantage with ML features; local partnerships; data moat compounds"),
    ("Food Safety Liability", "MODERATE", "Consumer illness from surplus food",
     "SFA license verification; merchant ToS; use-by date enforcement; food liability insurance recommended"),
    ("Low Consumer Trust", "MODERATE", "Consumers skeptical about surplus food quality",
     "Quality ratings; SFA verified badges; transparent pricing; risk-free first purchase"),
    ("Prediction Inaccuracy", "MODERATE", "Wrong surplus predictions frustrate merchants",
     "Confidence scores; easy override; prediction improves with data; rule-based fallback"),
    ("Grab/Foodpanda Feature", "LOW", "Established platforms add surplus food feature",
     "ML differentiation; merchant first focus; they optimize for delivery, not waste reduction"),
]

for i, (risk, level, desc, mitigation) in enumerate(risks):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 1.9)
    level_color = RED_ACCENT if level == "HIGH" else ORANGE if level == "MODERATE" else MED_GREEN

    add_rounded(sl, x, y, Inches(6.0), Inches(1.7), CARD_BG)
    tb(sl, x + Inches(0.1), y + Inches(0.05), Inches(4.5), Inches(0.3), risk, sz=13, color=DARK_GREEN, bold=True)
    add_rounded(sl, x + Inches(4.8), y + Inches(0.05), Inches(1.0), Inches(0.3), level_color)
    tb(sl, x + Inches(4.8), y + Inches(0.05), Inches(1.0), Inches(0.3), level, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, x + Inches(0.1), y + Inches(0.35), Inches(5.8), Inches(0.4), desc, sz=10, color=MED_TEXT)
    tb(sl, x + Inches(0.1), y + Inches(0.75), Inches(5.8), Inches(0.3), "Mitigation:", sz=10, color=MED_GREEN, bold=True)
    tb(sl, x + Inches(0.1), y + Inches(1.0), Inches(5.8), Inches(0.5), mitigation, sz=10, color=DARK_TEXT)


# ============================================================
# SLIDE 16: CLOSING / DEMO VISION
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_shape(sl, Inches(0), Inches(0), Inches(13.333), Inches(7.5), fill_color=DARK_GREEN)

tb(sl, Inches(1), Inches(0.8), Inches(11.3), Inches(0.6),
   "Demo Vision & Key Takeaways", sz=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# Demo scenario
add_rounded(sl, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.0), RGBColor(0x14, 0x4D, 0x1A))
tb(sl, Inches(1.0), Inches(1.65), Inches(11.3), Inches(0.35), "What the Prototype Will Demonstrate", sz=16, color=ACCENT_GREEN, bold=True)
tb(sl, Inches(1.0), Inches(2.1), Inches(11.3), Inches(1.3),
   "A working marketplace with synthetic data simulating 6 months of operation:\n"
   "• Merchant receives ML surplus prediction → approves listing → consumer discovers deal via personalized feed\n"
   "• Dynamic pricing adjusts in real-time as time passes and demand signals change\n"
   "• Waste analytics dashboard shows prescriptive recommendations with estimated savings",
   sz=13, color=RGBColor(0xC8, 0xE6, 0xC9))

# Key takeaways
takeaways = [
    ("Real Problem", "755K tonnes of food waste in Singapore, growing regulatory pressure"),
    ("Clear Gap", "No ML-powered food waste platform exists in Singapore or globally"),
    ("ML Moat", "4 ML capabilities create compounding data advantage that's hard to replicate"),
    ("Viable Business", "Break-even achievable at 100 merchants under moderate assumptions"),
]
for i, (title, desc) in enumerate(takeaways):
    x = Inches(0.8 + (i % 2) * 6.0)
    y = Inches(3.9 + (i // 2) * 1.4)
    add_rounded(sl, x, y, Inches(5.7), Inches(1.2), RGBColor(0x14, 0x4D, 0x1A))
    tb(sl, x + Inches(0.15), y + Inches(0.05), Inches(5.4), Inches(0.35), title, sz=15, color=ACCENT_GREEN, bold=True)
    tb(sl, x + Inches(0.15), y + Inches(0.4), Inches(5.4), Inches(0.6), desc, sz=12, color=RGBColor(0xC8, 0xE6, 0xC9))

# Bottom
tb(sl, Inches(1), Inches(6.6), Inches(11.3), Inches(0.5),
   "Thank You  |  Questions?",
   sz=24, color=WHITE, bold=True, align=PP_ALIGN.CENTER)


# ── Save ──
output_path = os.path.expanduser(
    "~/Desktop/SMU MBA Machine Learning/Individual Project/Food App/"
    "workspaces/Food App/Food_Waste_Marketplace_Analysis.pptx"
)
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
